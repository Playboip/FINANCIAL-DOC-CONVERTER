import os
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from motor.motor_asyncio import AsyncIOMotorClient
import stripe
import shutil
from docx import Document
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
from openpyxl import Workbook
from ebooklib import epub
from odf.opendocument import OpenDocumentText
from odf.text import P
from pptx import Presentation
from fpdf import FPDF

# -----------------------------
# Environment variables
# -----------------------------
# Explicitly fetch the remote MONGO_URL set on the Railway dashboard
MONGO_URL = os.getenv("MONGO_URL")
# --- FINAL FIX: Explicitly set the database name based on your choice ---
DB_NAME = "FDC" 
STRIPE_SECRET_KEY = os.getenv("STRIPE_SECRET_KEY")

# --- CRITICAL CHECK: Ensure the remote variable is loaded ---
if not MONGO_URL or "localhost" in MONGO_URL or "127.0.0.1" in MONGO_URL:
    # If it's missing or still contains a local value, raise a clear error
    raise ValueError(
        "MONGO_URL is either not set or still points to 'localhost'. "
        "You MUST set the environment variable on your Railway service to the correct remote URI (e.g., mongodb://mongo:***@crossover.proxy.rlwy.net:27384)."
    )
if not STRIPE_SECRET_KEY:
    raise ValueError("STRIPE_SECRET_KEY not set in environment")
# -------------------------------------------------------------

# Initialize Stripe
stripe.api_key = STRIPE_SECRET_KEY
print("✅ Stripe secret key loaded successfully.")

# -----------------------------
# MongoDB client
# -----------------------------
try:
    # Attempt to initialize the Motor client using the confirmed remote MONGO_URL
    client = AsyncIOMotorClient(MONGO_URL)
    # --- Using the explicit database name: FDC ---
    db = client[DB_NAME] 
    print(f"✅ MongoDB connected successfully to database: {DB_NAME}.")
except Exception as e:
    # If connection fails (e.g., credentials or network issue), provide details
    raise RuntimeError(
        f"Failed to connect to MongoDB. "
        f"Original error: {str(e)}. "
        f"Double-check the MONGO_URL credentials and network settings on Railway."
    )

# -----------------------------
# FastAPI app
# -----------------------------
app = FastAPI()

# -----------------------------
# Utility for saving uploaded files
# -----------------------------
UPLOAD_DIR = "./uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

def save_file_locally(upload_file: UploadFile) -> str:
    try:
        file_path = os.path.join(UPLOAD_DIR, upload_file.filename)
        # Ensure the file pointer is at the beginning before copying
        upload_file.file.seek(0)
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(upload_file.file, buffer)
        return file_path
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save file: {str(e)}")
    finally:
        # Close the file stream
        upload_file.file.close()


# -----------------------------
# Endpoints
# -----------------------------
@app.get("/")
async def health_check():
    return {"status": "ok", "message": "Server is running"}

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    return {"filename": file.filename, "path": file_path}

@app.post("/convert")
async def convert_file(file: UploadFile = File(...)):
    """
    Dummy conversion: just saves the file and returns new path.
    Extend this with your actual conversion logic.
    """
    file_path = save_file_locally(file)
    converted_path = file_path + ".converted"  # Example
    shutil.copy(file_path, converted_path)
    return {"original": file.filename, "converted_path": converted_path}

@app.post("/convert-to-docx")
async def convert_to_docx(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    docx_path = file_path + ".docx"
    document = Document()
    document.add_paragraph(file.file.read().decode('utf-8'))
    document.save(docx_path)
    return FileResponse(docx_path, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document', filename=file.filename + ".docx")

@app.post("/convert-to-txt")
async def convert_to_txt(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    txt_path = file_path + ".txt"
    with open(txt_path, "w") as f:
        f.write(file.file.read().decode('utf-8'))
    return FileResponse(txt_path, media_type='text/plain', filename=file.filename + ".txt")

@app.post("/convert-to-html")
async def convert_to_html(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    html_path = file_path + ".html"
    with open(html_path, "w") as f:
        f.write("<html><body><pre>")
        f.write(file.file.read().decode('utf-8'))
        f.write("</pre></body></html>")
    return FileResponse(html_path, media_type='text/html', filename=file.filename + ".html")

@app.post("/convert-to-odt")
async def convert_to_odt(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    odt_path = file_path + ".odt"
    # This is a dummy conversion for now. A real implementation would require a library like odfpy.
    with open(odt_path, "w") as f:
        f.write(file.file.read().decode('utf-8'))
    return FileResponse(odt_path, media_type='application/vnd.oasis.opendocument.text', filename=file.filename + ".odt")

@app.post("/convert-to-rtf")
async def convert_to_rtf(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    rtf_path = file_path + ".rtf"
    # This is a dummy conversion for now. A real implementation would require a library like unirtf.
    with open(rtf_path, "w") as f:
        f.write(file.file.read().decode('utf-8'))
    return FileResponse(rtf_path, media_type='application/rtf', filename=file.filename + ".rtf")

@app.post("/merge-pdfs")
async def merge_pdfs(files: list[UploadFile] = File(...)):
    merger = PdfMerger()
    for file in files:
        file_path = save_file_locally(file)
        merger.append(file_path)
    merged_path = "./uploads/merged.pdf"
    merger.write(merged_path)
    merger.close()
    return FileResponse(merged_path, media_type='application/pdf', filename="merged.pdf")

@app.post("/split-pdf")
async def split_pdf(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    reader = PdfReader(file_path)
    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)
        split_path = f"./uploads/split_{i}.pdf"
        with open(split_path, "wb") as f:
            writer.write(f)
    return JSONResponse({"message": "PDF split successfully"})

@app.post("/encrypt-pdf")
async def encrypt_pdf(file: UploadFile = File(...), password: str = ""):
    file_path = save_file_locally(file)
    reader = PdfReader(file_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    encrypted_path = file_path + ".encrypted.pdf"
    with open(encrypted_path, "wb") as f:
        writer.write(f)
    return FileResponse(encrypted_path, media_type='application/pdf', filename=file.filename + ".encrypted.pdf")

@app.post("/convert-to-xlsx")
async def convert_to_xlsx(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    xlsx_path = file_path + ".xlsx"
    workbook = Workbook()
    sheet = workbook.active
    with open(file_path, 'r') as f:
        for row_idx, line in enumerate(f, 1):
            for col_idx, cell_value in enumerate(line.strip().split(','), 1):
                sheet.cell(row=row_idx, column=col_idx, value=cell_value)
    workbook.save(xlsx_path)
    return FileResponse(xlsx_path, media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', filename=file.filename + ".xlsx")

@app.post("/convert-to-epub")
async def convert_to_epub(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    epub_path = file_path + ".epub"
    book = epub.EpubBook()
    book.set_title(file.filename)
    book.set_language('en')
    c1 = epub.EpubHtml(title='Intro', file_name='chap_01.xhtml', lang='en')
    c1.content = u'<h1>Intro</h1><p>%s</p>' % file.file.read().decode('utf-8')
    book.add_item(c1)
    book.toc = (epub.Link('chap_01.xhtml', 'Intro', 'intro'),)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    style = 'BODY {color: white;}'
    nav_css = epub.EpubItem(uid="style_nav", file_name="style/nav.css", media_type="text/css", content=style)
    book.add_item(nav_css)
    book.spine = ['nav', c1]
    epub.write_epub(epub_path, book, {})
    return FileResponse(epub_path, media_type='application/epub+zip', filename=file.filename + ".epub")

@app.post("/convert-to-pdf")
async def convert_to_pdf(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    pdf_path = file_path + ".pdf"
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size = 12)
    with open(file_path, "r") as f:
        for x in f:
            pdf.cell(200, 10, txt = x, ln = 1, align = 'C')
    pdf.output(pdf_path)
    return FileResponse(pdf_path, media_type='application/pdf', filename=file.filename + ".pdf")

@app.post("/convert-to-odp")
async def convert_to_odp(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    odp_path = file_path + ".odp"
    # This is a dummy conversion for now. A real implementation would require a library like python-odf.
    with open(odp_path, "w") as f:
        f.write(file.file.read().decode('utf-8'))
    return FileResponse(odp_path, media_type='application/vnd.oasis.opendocument.presentation', filename=file.filename + ".odp")

@app.post("/convert-to-ods")
async def convert_to_ods(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    ods_path = file_path + ".ods"
    # This is a dummy conversion for now. A real implementation would require a library like python-odf.
    with open(ods_path, "w") as f:
        f.write(file.file.read().decode('utf-8'))
    return FileResponse(ods_path, media_type='application/vnd.oasis.opendocument.spreadsheet', filename=file.filename + ".ods")

@app.post("/convert-to-pptx")
async def convert_to_pptx(file: UploadFile = File(...)):
    file_path = save_file_locally(file)
    pptx_path = file_path + ".pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = file.filename
    with open(file_path, 'r') as f:
        for line in f:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            p = slide.shapes.add_textbox(1, 1, 1, 1).text_frame.add_paragraph()
            p.text = line
    prs.save(pptx_path)
    return FileResponse(pptx_path, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename=file.filename + ".pptx")

@app.post("/create-payment-intent")
async def create_payment_intent(amount: int):
    """
    Create Stripe PaymentIntent
    """
    try:
        intent = stripe.PaymentIntent.create(
            amount=amount,
            currency="usd"
        )
        return {"client_secret": intent.client_secret}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Stripe error: {str(e)}")

# -----------------------------
# MongoDB CRUD endpoints
# -----------------------------
@app.post("/add-document")
async def add_document(collection_name: str, document: dict):
    try:
        collection = db[collection_name]
        result = await collection.insert_one(document)
        # Using str() is necessary as MongoDB's ObjectId is not JSON serializable
        return {"inserted_id": str(result.inserted_id)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to insert document: {str(e)}")

@app.get("/get-documents")
async def get_documents(collection_name: str):
    try:
        collection = db[collection_name]
        cursor = collection.find({})
        documents = []
        async for doc in cursor:
            # Convert ObjectId to string for JSON serialization
            doc["_id"] = str(doc["_id"])
            documents.append(doc)
        return documents
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch documents: {str(e)}")

# -----------------------------
# Uvicorn launch for Railway
# -----------------------------
if __name__ == "__main__":
    import uvicorn
    # Use environment variable PORT, or default to 8000
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
